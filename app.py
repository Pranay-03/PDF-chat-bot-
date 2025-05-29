# Installation requirements:
# pip install streamlit streamlit-chat streamlit-webrtc
# pip install pymupdf python-docx python-pptx pandas openpyxl xlrd striprtf beautifulsoup4
# pip install google-generativeai pinecone-client tiktoken scikit-learn numpy
# pip install pydub speechrecognition gtts

# main.py - Complete RAG Chatbot with Gemini & WebRTC
import streamlit as st
import sqlite3
import json
from datetime import datetime,date
import google.generativeai as genai
from pinecone.grpc import PineconeGRPC as Pinecone
from pinecone import ServerlessSpec
import time
from typing import List, Dict, Any
import uuid
import tiktoken
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import speech_recognition as sr
from gtts import gTTS
import io
import tempfile
import os

# WebRTC imports
from streamlit_webrtc import webrtc_streamer, WebRtcMode, RTCConfiguration
import av
import pydub
from pydub import AudioSegment

# Document processors
import fitz  # pymupdf
from docx import Document
import pandas as pd
from pptx import Presentation

class AccessibleRAGChatbot:
    def __init__(self):
        self.setup_database()
        self.setup_apis()
        self.setup_accessibility()
    
    def setup_database(self):
        """Initialize SQLite for chat history"""
        # --- FIX: DeprecationWarning for datetime adapter in sqlite3 ---
        sqlite3.register_adapter(datetime, lambda ts: ts.isoformat())
        sqlite3.register_adapter(date, lambda d: d.isoformat()) # Also register date if you use date objects
        
        def convert_timestamp(val):
            return datetime.fromisoformat(val.decode('utf-8'))
        
        sqlite3.converters['DATETIME'] = convert_timestamp
        sqlite3.converters['TIMESTAMP'] = convert_timestamp # Add if you also use TIMESTAMP columns

        self.conn = sqlite3.connect('chat_history.db', check_same_thread=False,
                                    detect_types=sqlite3.PARSE_DECLTYPES | sqlite3.PARSE_COLNAMES)
        # --- END FIX ---

        self.conn.execute('''
            CREATE TABLE IF NOT EXISTS conversations (
                id TEXT PRIMARY KEY,
                timestamp DATETIME,
                user_message TEXT,
                bot_response TEXT,
                context_docs TEXT,
                session_id TEXT
            )
        ''')
        self.conn.commit()
    

    
    def setup_apis(self):
        """Initialize API connections"""
        # Configure Gemini API
        if 'GEMINI_API_KEY' in st.secrets:
            genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
            st.success("✅ Gemini API configured")
        else:
            st.error("❌ Gemini API key not found in secrets")
        
        # Initialize components
        if 'vector_store' not in st.session_state:
            st.session_state.vector_store = VectorStore()
        if 'llm_handler' not in st.session_state:
            st.session_state.llm_handler = LLMHandler()
        if 'chunker' not in st.session_state:
            st.session_state.chunker = DocumentChunker()
    
    def setup_accessibility(self):
        """Configure accessibility settings"""
        if 'accessibility_settings' not in st.session_state:
            st.session_state.accessibility_settings = {
                'high_contrast': False,
                'large_text': False,
                'screen_reader_mode': False,
                'audio_speed': 1.0,
                'auto_play_responses': True
            }

class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(file) -> str:
        """Extract text from PDF using pymupdf"""
        try:
            doc = fitz.open(stream=file.read(), filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            return text
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error processing Word document: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_excel(file) -> str:
        """Extract text from Excel file"""
        try:
            excel_file = pd.ExcelFile(file)
            text = ""
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file, sheet_name=sheet_name)
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string() + "\n\n"
            return text
        except Exception as e:
            st.error(f"Error processing Excel file: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file) -> str:
        """Extract text from PowerPoint"""
        try:
            prs = Presentation(file)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except Exception as e:
            st.error(f"Error processing PowerPoint: {str(e)}")
            return ""
    
    @staticmethod
    def process_document(file) -> str:
        """Route document to appropriate processor"""
        file_extension = file.name.split('.')[-1].lower()
        
        processors = {
            'pdf': DocumentProcessor.extract_text_from_pdf,
            'docx': DocumentProcessor.extract_text_from_docx,
            'doc': DocumentProcessor.extract_text_from_docx,
            'xlsx': DocumentProcessor.extract_text_from_excel,
            'xls': DocumentProcessor.extract_text_from_excel,
            'pptx': DocumentProcessor.extract_text_from_pptx,
            'txt': lambda f: f.read().decode('utf-8')
        }
        
        if file_extension in processors:
            return processors[file_extension](file)
        else:
            st.error(f"Unsupported file type: {file_extension}")
            return ""

class DocumentChunker:
    def __init__(self, chunk_size: int = 1000, overlap: int = 200):
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
    
    def count_tokens(self, text: str) -> int:
        """Count tokens in text"""
        return len(self.tokenizer.encode(text))
    
    def chunk_text(self, text: str, source_name: str) -> List[Dict[str, Any]]:
        """Chunk text into overlapping segments"""
        if not text.strip():
            return []
        
        sentences = text.replace('\n', ' ').split('. ')
        chunks = []
        current_chunk = ""
        current_tokens = 0
        
        for sentence in sentences:
            sentence_tokens = self.count_tokens(sentence)
            
            if current_tokens + sentence_tokens > self.chunk_size and current_chunk:
                chunks.append({
                    'text': current_chunk.strip(),
                    'source': source_name,
                    'tokens': current_tokens,
                    'chunk_id': f"{source_name}_{len(chunks)}"
                })
                
                overlap_text = self._get_overlap_text(current_chunk)
                current_chunk = overlap_text + sentence + ". "
                current_tokens = self.count_tokens(current_chunk)
            else:
                current_chunk += sentence + ". "
                current_tokens += sentence_tokens
        
        if current_chunk.strip():
            chunks.append({
                'text': current_chunk.strip(),
                'source': source_name,
                'tokens': current_tokens,
                'chunk_id': f"{source_name}_{len(chunks)}"
            })
        
        return chunks
    
    def _get_overlap_text(self, text: str) -> str:
        """Get overlap text from end of chunk"""
        words = text.split()
        overlap_words = words[-int(len(words) * 0.2):]
        return " ".join(overlap_words) + " "

class VectorStore:
    def __init__(self):
        self.pc = None
        self.index = None
        self.local_vectors = {}
        self.setup_pinecone()
    
    def setup_pinecone(self):
        """Initialize Pinecone vector database with new API"""
        try:
            if 'PINECONE_API_KEY' in st.secrets:
                # Initialize Pinecone client
                self.pc = Pinecone(api_key=st.secrets["PINECONE_API_KEY"])
                
                index_name = "rag-chatbot-gemini"
                
                # Check if index exists
                existing_indexes = [index.name for index in self.pc.list_indexes()]
                
                if index_name not in existing_indexes:
                    # Create index with new API
                    self.pc.create_index(
                        name=index_name,
                        dimension=768,  # Gemini embedding dimension
                        metric="cosine",
                        spec=ServerlessSpec(
                            cloud="aws",
                            region="us-east-1"
                        )
                    )
                    # Wait for index to be ready
                    time.sleep(10)
                
                # Connect to index
                self.index = self.pc.Index(index_name)
                st.success("✅ Connected to Pinecone vector database")
                
        except Exception as e:
            st.warning(f"Pinecone setup failed, using local storage: {str(e)}")
            self.pc = None
            self.index = None
    
    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate embeddings using Gemini API"""
        try:
            embeddings = []
            for text in texts:
                # Use Gemini's embedding model
                result = genai.embed_content(
                    model="models/text-embedding-004",
                    content=text,
                    task_type="retrieval_document"
                )
                embeddings.append(result['embedding'])
            return embeddings
        except Exception as e:
            st.error(f"Gemini embedding generation failed: {str(e)}")
            return []
    
    def store_chunks(self, chunks: List[Dict[str, Any]]) -> bool:
        """Store document chunks with embeddings"""
        if not chunks:
            return False
        
        texts = [chunk['text'] for chunk in chunks]
        embeddings = self.generate_embeddings(texts)
        
        if not embeddings:
            return False
        
        # Store in Pinecone if available
        if self.index:
            try:
                vectors_to_upsert = []
                for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                    vectors_to_upsert.append({
                        'id': chunk['chunk_id'],
                        'values': embedding,
                        'metadata': {
                            'text': chunk['text'],
                            'source': chunk['source'],
                            'tokens': chunk['tokens']
                        }
                    })
                
                # Upsert in batches
                batch_size = 100
                for i in range(0, len(vectors_to_upsert), batch_size):
                    batch = vectors_to_upsert[i:i + batch_size]
                    self.index.upsert(vectors=batch)
                
                return True
            except Exception as e:
                st.error(f"Pinecone storage failed: {str(e)}")
        
        # Fallback to local storage
        for chunk, embedding in zip(chunks, embeddings):
            self.local_vectors[chunk['chunk_id']] = {
                'embedding': embedding,
                'metadata': {
                    'text': chunk['text'],
                    'source': chunk['source'],
                    'tokens': chunk['tokens']
                }
            }
        
        return True
    
    def similarity_search(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Search for similar chunks"""
        try:
            query_embedding = genai.embed_content(
                model="models/text-embedding-004",
                content=query,
                task_type="retrieval_query"
            )['embedding']
        except Exception as e:
            st.error(f"Query embedding failed: {str(e)}")
            return []
        
        # Search in Pinecone if available
        if self.index:
            try:
                results = self.index.query(
                    vector=query_embedding,
                    top_k=top_k,
                    include_metadata=True
                )
                
                return [{
                    'text': match['metadata']['text'],
                    'source': match['metadata']['source'],
                    'score': match['score'],
                    'tokens': match['metadata'].get('tokens', 0)
                } for match in results['matches']]
            except Exception as e:
                st.error(f"Pinecone search failed: {str(e)}")
        
        # Fallback to local search
        if not self.local_vectors:
            return []
        
        similarities = []
        for chunk_id, data in self.local_vectors.items():
            similarity = cosine_similarity(
                [query_embedding], 
                [data['embedding']]
            )[0][0]
            
            similarities.append({
                'text': data['metadata']['text'],
                'source': data['metadata']['source'],
                'score': similarity,
                'tokens': data['metadata']['tokens'],
                'chunk_id': chunk_id
            })
        
        similarities.sort(key=lambda x: x['score'], reverse=True)
        return similarities[:top_k]

class LLMHandler:
    def __init__(self):
        self.model = genai.GenerativeModel('gemini-2.5-flash-preview-05-20')
    
    def generate_response(self, query: str, context_chunks: List[Dict[str, Any]], 
                         conversation_history: List[Dict[str, str]]) -> str:
        """Generate response using Gemini 2.5 Flash"""
        
        # Build context from retrieved chunks
        context_text = "\n\n".join([
            f"Source: {chunk['source']}\nContent: {chunk['text']}"
            for chunk in context_chunks
        ])
        
        # Build conversation context
        conversation_context = ""
        if conversation_history:
            conversation_context = "Previous conversation:\n"
            for msg in conversation_history[-3:]:
                conversation_context += f"User: {msg['user']}\nAssistant: {msg['bot']}\n\n"
        
        # Create comprehensive prompt
        prompt = f"""You are a helpful AI assistant that answers questions based on provided documents and conversation context.

Guidelines:
1. Use the provided document context to answer questions accurately
2. Reference specific sources when making claims
3. If the information isn't in the documents, say so clearly
4. Consider the conversation history for context
5. Be conversational and helpful
6. If asked about previous topics, refer to the conversation history

{conversation_context}

Document Context:
{context_text}

Current Question: {query}

Please provide a comprehensive answer based on the document context and conversation history."""
        
        try:
            response = self.model.generate_content(prompt)
            return response.text
        except Exception as e:
            st.error(f"Gemini API error: {str(e)}")
            return "I apologize, but I'm having trouble generating a response. Please try again."

class VoiceHandler:
    def __init__(self):
        self.recognizer = sr.Recognizer()
    
    def process_audio_frame(self, audio_frame):
        """Process audio frame from WebRTC"""
        try:
            # Convert audio frame to wav
            audio_segment = AudioSegment(
                audio_frame.to_ndarray().tobytes(),
                frame_rate=audio_frame.sample_rate,
                sample_width=audio_frame.format.bytes,
                channels=len(audio_frame.layout.channels)
            )
            
            # Convert to recognizable format
            wav_io = io.BytesIO()
            audio_segment.export(wav_io, format="wav")
            wav_io.seek(0)
            
            # Recognize speech
            with sr.AudioFile(wav_io) as source:
                audio_data = self.recognizer.record(source)
                text = self.recognizer.recognize_google(audio_data)
                return text
        except Exception as e:
            return f"Error processing audio: {str(e)}"
    
    def text_to_speech(self, text: str, lang: str = "en") -> bytes:
        """Convert text to speech using gTTS"""
        try:
            tts = gTTS(text=text, lang=lang, slow=False)
            audio_buffer = io.BytesIO()
            tts.write_to_fp(audio_buffer)
            audio_buffer.seek(0)
            return audio_buffer.read()
        except Exception as e:
            st.error(f"Text-to-speech error: {str(e)}")
            return b""

class ContextManager:
    def __init__(self, db_connection):
        self.conn = db_connection
    
    def save_conversation(self, user_message: str, bot_response: str, 
                         context_docs: List[str], session_id: str):
        """Save conversation to database"""
        conversation_id = str(uuid.uuid4())
        self.conn.execute('''
            INSERT INTO conversations 
            (id, timestamp, user_message, bot_response, context_docs, session_id)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (conversation_id, datetime.now(), user_message, bot_response, 
              json.dumps(context_docs), session_id))
        self.conn.commit()
    
    def get_conversation_history(self, session_id: str, limit: int = 10) -> List[Dict]:
        """Retrieve conversation history"""
        cursor = self.conn.execute('''
            SELECT user_message, bot_response, timestamp 
            FROM conversations 
            WHERE session_id = ? 
            ORDER BY timestamp DESC 
            LIMIT ?
        ''', (session_id, limit))
        
        return [{"user": row[0], "bot": row[1], "timestamp": row[2]} 
                for row in cursor.fetchall()]

# WebRTC Audio Processor
class AudioProcessor:
    def __init__(self):
        self.voice_handler = VoiceHandler()
        self.audio_buffer = []
        
    def recv(self, frame):
        """Process incoming audio frame"""
        try:
            # Store audio frame
            self.audio_buffer.append(frame)
            
            # Process if buffer has enough frames
            if len(self.audio_buffer) >= 10:  # Process every 10 frames
                # Combine frames and process
                combined_audio = self._combine_frames(self.audio_buffer)
                text = self.voice_handler.process_audio_frame(combined_audio)
                
                # Store recognized text
                if text and 'webrtc_text' not in st.session_state:
                    st.session_state.webrtc_text = text
                
                # Clear buffer
                self.audio_buffer = []
        except Exception as e:
            st.error(f"Audio processing error: {str(e)}")
        
        return frame
    
    def _combine_frames(self, frames):
        """Combine multiple audio frames"""
        # Simple implementation - return first frame for now
        return frames[0] if frames else None

def main():
    # Page config with accessibility
    st.set_page_config(
        page_title="Gemini RAG Chatbot with WebRTC",
        page_icon="🎤",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Initialize components
    if 'chatbot' not in st.session_state:
        st.session_state.chatbot = AccessibleRAGChatbot()
        st.session_state.doc_processor = DocumentProcessor()
        st.session_state.voice_handler = VoiceHandler()
        st.session_state.context_manager = ContextManager(st.session_state.chatbot.conn)
        st.session_state.session_id = str(uuid.uuid4())
        st.session_state.messages = []
        st.session_state.documents = {}
        st.session_state.webrtc_text = ""
    
    # Accessibility Controls Sidebar
    with st.sidebar:
        st.header("♿ Accessibility Settings")
        
        # Visual accessibility
        st.session_state.accessibility_settings['high_contrast'] = st.checkbox(
            "High Contrast Mode", 
            value=st.session_state.accessibility_settings['high_contrast']
        )
        
        st.session_state.accessibility_settings['large_text'] = st.checkbox(
            "Large Text Mode",
            value=st.session_state.accessibility_settings['large_text']
        )
        
        # Audio accessibility
        st.session_state.accessibility_settings['auto_play_responses'] = st.checkbox(
            "Auto-play Audio Responses",
            value=st.session_state.accessibility_settings['auto_play_responses']
        )
        
        st.divider()
        
        # Document Upload
        st.header("📁 Upload Documents")
        uploaded_files = st.file_uploader(
            "Choose files",
            accept_multiple_files=True,
            type=['pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'txt'],
            help="Supported formats: PDF, Word, Excel, PowerPoint, Text"
        )
        
        if uploaded_files:
            for file in uploaded_files:
                if file.name not in st.session_state.documents:
                    with st.spinner(f"Processing {file.name}..."):
                        text_content = st.session_state.doc_processor.process_document(file)
                        if text_content:
                            st.session_state.documents[file.name] = text_content
                            
                            with st.spinner(f"Creating Gemini embeddings for {file.name}..."):
                                chunks = st.session_state.chunker.chunk_text(text_content, file.name)
                                if chunks:
                                    success = st.session_state.vector_store.store_chunks(chunks)
                                    if success:
                                        st.success(f"✅ {file.name} processed and indexed! ({len(chunks)} chunks)")
                                    else:
                                        st.error(f"❌ Failed to index {file.name}")
                                else:
                                    st.warning(f"⚠️ No chunks created for {file.name}")
        
        # Show processed documents
        if st.session_state.documents:
            st.subheader("📚 Processed Documents")
            for doc_name in st.session_state.documents.keys():
                st.text(f"• {doc_name}")
    
    # Apply accessibility styles
    if st.session_state.accessibility_settings['high_contrast']:
        st.markdown("""
        <style>
        .stApp { background-color: black; color: white; }
        .stTextInput > div > div > input { background-color: #333; color: white; }
        </style>
        """, unsafe_allow_html=True)
    
    if st.session_state.accessibility_settings['large_text']:
        st.markdown("""
        <style>
        html, body, [class*="css"] { font-size: 18px !important; }
        </style>
        """, unsafe_allow_html=True)
    
    # Main Interface
    st.title("🎤 Gemini RAG Chatbot with WebRTC")
    st.caption("Voice-enabled document Q&A powered by Gemini 2.0 Flash")
    
    # Voice Input Section
    st.header("🎙️ Voice Input")
    
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.subheader("🌐 WebRTC Live Audio")
        
        # WebRTC Audio Streamer
        webrtc_ctx = webrtc_streamer(
            key="speech-to-text",
            mode=WebRtcMode.SENDONLY,
            audio_processor_factory=AudioProcessor,
            rtc_configuration=RTCConfiguration({
                "iceServers": [{"urls": ["stun:stun.l.google.com:19302"]}]
            }),
            media_stream_constraints={"video": False, "audio": True},
            async_processing=True,
        )
        
        # Display recognized text from WebRTC
        if st.session_state.webrtc_text:
            st.success(f"🎯 Recognized: {st.session_state.webrtc_text}")
            
            if st.button("📝 Use WebRTC Text", use_container_width=True):
                st.session_state.messages.append({
                    "role": "user", 
                    "content": st.session_state.webrtc_text
                })
                st.session_state.webrtc_text = ""  # Clear after use
                st.rerun()
    
    with col2:
        st.subheader("📝 Text Input")
        
        # Text input alternative
        text_input = st.text_input(
            "Type your question:", 
            key="text_question",
            help="You can type your question here as an alternative to voice input"
        )
        
        if text_input and st.button("🚀 Send Message", use_container_width=True):
            st.session_state.messages.append({"role": "user", "content": text_input})
            st.rerun()
    
    # Chat History Display
    st.header("💬 Conversation")
    
    # Display chat messages
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.write(message["content"])
            
            # Add audio playback for bot responses
            if (message["role"] == "assistant" and 
                st.session_state.accessibility_settings['auto_play_responses']):
                audio_content = st.session_state.voice_handler.text_to_speech(message["content"])
                if audio_content:
                    st.audio(audio_content, format="audio/mp3")
    
    # Process new messages
    if st.session_state.messages and st.session_state.messages[-1]["role"] == "user":
        user_question = st.session_state.messages[-1]["content"]
        
        with st.chat_message("assistant"):
            with st.spinner("🔍 Searching relevant documents with Gemini embeddings..."):
                # Search for relevant chunks
                relevant_chunks = st.session_state.vector_store.similarity_search(
                    user_question, top_k=5
                )
                
                if relevant_chunks:
                    st.info(f"📚 Found {len(relevant_chunks)} relevant chunks from your documents")
                    sources = list(set([chunk['source'] for chunk in relevant_chunks]))
                    st.caption(f"Sources: {', '.join(sources)}")
                else:
                    st.warning("⚠️ No relevant content found in your documents")
                
                # Get conversation history
                history = st.session_state.context_manager.get_conversation_history(
                    st.session_state.session_id
                )
                
                with st.spinner("🤖 Generating response with Gemini 2.0 Flash..."):
                    # Generate response using Gemini
                    bot_response = st.session_state.llm_handler.generate_response(
                        user_question, 
                        relevant_chunks, 
                        history
                    )
                    
                    st.write(bot_response)
                    
                    # Save conversation
                    st.session_state.context_manager.save_conversation(
                        user_question, bot_response, 
                        [chunk['source'] for chunk in relevant_chunks],
                        st.session_state.session_id
                    )
                    
                    # Add response to messages
                    st.session_state.messages.append({"role": "assistant", "content": bot_response})
                    
                    # Generate and play audio response
                    if st.session_state.accessibility_settings['auto_play_responses']:
                        with st.spinner("🔊 Converting to speech..."):
                            audio_content = st.session_state.voice_handler.text_to_speech(bot_response)
                            if audio_content:
                                st.audio(audio_content, format="audio/mp3")
                
                # Show relevant chunks in expander
                if relevant_chunks:
                    with st.expander("📖 View Source Context"):
                        for i, chunk in enumerate(relevant_chunks):
                            st.markdown(f"""
                            **Source {i+1}: {chunk['source']}** (Similarity: {chunk['score']:.3f})
                            
                            {chunk['text'][:500]}{'...' if len(chunk['text']) > 500 else ''}
                            
                            ---
                            """)
    
    # Clear Chat Button
    if st.button("🗑️ Clear Chat", use_container_width=True):
        st.session_state.messages = []
        st.session_state.webrtc_text = ""
        st.rerun()
    
    # Info Section
    with st.expander("ℹ️ About This App"):
        st.markdown("""
        **🚀 Features:**
        - **Gemini 2.5 Flash** for intelligent responses
        - **Gemini Text Embedding** for document search
        - **WebRTC** for real-time voice input
        - **Multi-format** document support (PDF, Word, Excel, PowerPoint)
        - **Accessibility** features with audio feedback
        - **Context-aware** conversations with memory
        
        **🎤 How to Use:**
        1. Upload your documents using the sidebar
        2. Use WebRTC for real-time voice input or type your questions
        3. Get intelligent responses based on your documents
        4. Enjoy audio feedback for accessibility
        
        **🔧 Powered by:**
        - Google Gemini 2.0 Flash (LLM)
        - Google Text Embedding (Vector Search)
        - Streamlit WebRTC (Voice Input)
        - Pinecone (Vector Database)
        """)

if __name__ == "__main__":
    main()