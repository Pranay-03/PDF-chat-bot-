# Document processors
import fitz  # pymupdf
from docx import Document
import pandas as pd
from pptx import Presentation

from typing import List, Dict, Any
import streamlit as st
import tiktoken


class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(file) -> str:
        """Extract text from PDF using pymupdf"""
        try:
            doc = fitz.open(stream=file.read(), filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            return text
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error processing Word document: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_excel(file) -> str:
        """Extract text from Excel file"""
        try:
            excel_file = pd.ExcelFile(file)
            text = ""
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file, sheet_name=sheet_name)
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string() + "\n\n"
            return text
        except Exception as e:
            st.error(f"Error processing Excel file: {str(e)}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file) -> str:
        """Extract text from PowerPoint"""
        try:
            prs = Presentation(file)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except Exception as e:
            st.error(f"Error processing PowerPoint: {str(e)}")
            return ""
    
    @staticmethod
    def process_document(file) -> str:
        """Route document to appropriate processor"""
        file_extension = file.name.split('.')[-1].lower()
        
        processors = {
            'pdf': DocumentProcessor.extract_text_from_pdf,
            'docx': DocumentProcessor.extract_text_from_docx,
            'doc': DocumentProcessor.extract_text_from_docx,
            'xlsx': DocumentProcessor.extract_text_from_excel,
            'xls': DocumentProcessor.extract_text_from_excel,
            'pptx': DocumentProcessor.extract_text_from_pptx,
            'txt': lambda f: f.read().decode('utf-8')
        }
        
        if file_extension in processors:
            return processors[file_extension](file)
        else:
            st.error(f"Unsupported file type: {file_extension}")
            return ""

class DocumentChunker:
    def __init__(self, chunk_size: int = 1000, overlap: int = 200):
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
    
    def count_tokens(self, text: str) -> int:
        """Count tokens in text"""
        return len(self.tokenizer.encode(text))
    
    def chunk_text(self, text: str, source_name: str) -> List[Dict[str, Any]]:
        """Chunk text into overlapping segments"""
        if not text.strip():
            return []
        
        sentences = text.replace('\n', ' ').split('. ')
        chunks = []
        current_chunk = ""
        current_tokens = 0
        
        for sentence in sentences:
            sentence_tokens = self.count_tokens(sentence)
            
            if current_tokens + sentence_tokens > self.chunk_size and current_chunk:
                chunks.append({
                    'text': current_chunk.strip(),
                    'source': source_name,
                    'tokens': current_tokens,
                    'chunk_id': f"{source_name}_{len(chunks)}"
                })
                
                overlap_text = self._get_overlap_text(current_chunk)
                current_chunk = overlap_text + sentence + ". "
                current_tokens = self.count_tokens(current_chunk)
            else:
                current_chunk += sentence + ". "
                current_tokens += sentence_tokens
        
        if current_chunk.strip():
            chunks.append({
                'text': current_chunk.strip(),
                'source': source_name,
                'tokens': current_tokens,
                'chunk_id': f"{source_name}_{len(chunks)}"
            })
        
        return chunks
    
    def _get_overlap_text(self, text: str) -> str:
        """Get overlap text from end of chunk"""
        words = text.split()
        overlap_words = words[-int(len(words) * 0.2):]
        return " ".join(overlap_words) + " "